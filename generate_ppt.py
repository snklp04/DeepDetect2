"""
Generate Minor Project Evaluation PPT for DeepDetect2.
All content sourced from the actual codebase & project report.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7B, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MUTED_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_CARD = RGBColor(0x24, 0x24, 0x3E)
GREEN = RGBColor(0x00, 0xE6, 0x76)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)
RED_ACCENT = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW = RGBColor(0xFF, 0xD9, 0x3D)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=16,
                       default_color=LIGHT_GRAY, line_spacing=1.3):
    """lines: list of (text, font_size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Segoe UI"
        p.space_after = Pt(size * 0.4)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT_CYAN, spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(font_size * 0.5)
        run = p.add_run()
        run.text = ""
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_table_slide_content(slide, left, top, width, rows_data, col_widths, header_bg=ACCENT_BLUE):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.45 * n_rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw
    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.name = "Segoe UI"
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_CARD if ri % 2 == 1 else RGBColor(0x1E, 0x1E, 0x36)
    return table_shape


# ─── PRESENTATION ──────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

# ═══════════════════ SLIDE 1: TITLE ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_accent_line(slide, Inches(0), Inches(0), SLIDE_W, ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1),
             "DEEPFAKE DETECTION SYSTEM", 42, ACCENT_CYAN, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.8),
             "Multimodal Audio-Visual Deepfake Detection Framework",
             22, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5), Inches(3.2), Inches(3.333), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.5),
             "Minor Project Evaluation", 20, YELLOW, True, PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(2), [
    ("B.Tech in Computer Science and Engineering", 18, LIGHT_GRAY, False),
    ("", 8, LIGHT_GRAY, False),
    ("Sankalp (12311001)     |     Madhur Suman (12311002)", 18, WHITE, True),
    ("", 8, LIGHT_GRAY, False),
    ("Supervised by: Dr. Md. Arquam (Assistant Professor)", 16, MUTED_GRAY, False),
    ("", 8, LIGHT_GRAY, False),
    ("Indian Institute of Information Technology, Sonepat", 16, MUTED_GRAY, False),
    ("April 2026", 16, MUTED_GRAY, False),
])
for line in slide.shapes:
    if hasattr(line, 'text_frame'):
        for p in line.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER

# ═══════════════════ SLIDE 2: OUTLINE ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "PRESENTATION OUTLINE", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3), ACCENT_BLUE)

outline_items = [
    ("01", "Introduction & Problem Statement"),
    ("02", "Project Objectives"),
    ("03", "Architecture & Model Design"),
    ("04", "Self-Supervised Pretraining"),
    ("05", "Fine-Tuning & Classification"),
    ("06", "Results & Metrics"),
    ("07", "Tech Stack & Tools"),
    ("08", "Limitations & Future Scope"),
    ("09", "Conclusion"),
]

for i, (num, title) in enumerate(outline_items):
    y = Inches(1.7) + Inches(i * 0.58)
    add_rect(slide, Inches(1.2), y, Inches(10.5), Inches(0.48), DARK_CARD)
    add_text_box(slide, Inches(1.4), y + Pt(4), Inches(0.8), Inches(0.4),
                 num, 18, ACCENT_BLUE, True)
    add_text_box(slide, Inches(2.3), y + Pt(4), Inches(9), Inches(0.4),
                 title, 18, WHITE, False)

# ═══════════════════ SLIDE 3: PROBLEM STATEMENT ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "THE PROBLEM", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1),
             "Modern deepfakes are increasingly convincing visually, but often contain subtle audio-visual inconsistencies.",
             18, LIGHT_GRAY, False)

problems = [
    ("Lip Sync Errors", "Mouth movements don't match speech", RED_ACCENT),
    ("Audio Splicing", "Real audio + fake video misalignment", ORANGE),
    ("Temporal Glitches", "Unnatural frame transitions", YELLOW),
]
for i, (title, desc, color) in enumerate(problems):
    x = Inches(0.8) + Inches(i * 4)
    card = add_rect(slide, x, Inches(2.8), Inches(3.6), Inches(1.5), DARK_CARD, color)
    add_text_box(slide, x + Inches(0.2), Inches(2.95), Inches(3.2), Inches(0.5),
                 title, 18, color, True)
    add_text_box(slide, x + Inches(0.2), Inches(3.5), Inches(3.2), Inches(0.6),
                 desc, 14, LIGHT_GRAY, False)

challenges = [
    "Most existing detectors rely only on visual analysis and miss audio cues",
    "Lack of large, high-quality labelled deepfake datasets",
    "Visual-only methods fail to generalize to new generation techniques",
    "Computational cost limits real-world deployment",
]
add_text_box(slide, Inches(0.8), Inches(4.7), Inches(11), Inches(0.5),
             "Key Challenges:", 18, ACCENT_BLUE, True)
for i, ch in enumerate(challenges):
    add_text_box(slide, Inches(1.2), Inches(5.3) + Inches(i * 0.45), Inches(11), Inches(0.4),
                 f"\u2022  {ch}", 15, LIGHT_GRAY, False)

# ═══════════════════ SLIDE 4: OBJECTIVES ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "PROJECT OBJECTIVES", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3), ACCENT_BLUE)

objectives = [
    ("Cross-Modal Detection",
     "Analyse both video and audio streams together to catch subtle mismatches like lip sync errors and spliced audio using a Multimodal Bottleneck Transformer."),
    ("Self-Supervised Representation Learning",
     "Pretrain on unlabelled VoxCeleb videos using contrastive learning, synchronization prediction, and masked reconstruction to address data scarcity."),
    ("Practical Performance",
     "Achieve competitive accuracy, precision, recall, and F1 on FakeAVCeleb while keeping the system efficient for consumer GPUs (NVIDIA T4+)."),
    ("Modular Architecture",
     "Design clean, swappable components (visual encoder, audio encoder, fusion engine, classifier) so each can be upgraded independently."),
]

for i, (title, desc) in enumerate(objectives):
    y = Inches(1.6) + Inches(i * 1.35)
    card = add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(1.15), DARK_CARD, ACCENT_BLUE)
    add_text_box(slide, Inches(1.1), y + Pt(6), Inches(11), Inches(0.4),
                 f"{i+1}.  {title}", 18, ACCENT_CYAN, True)
    add_text_box(slide, Inches(1.4), y + Pt(30), Inches(10.8), Inches(0.7),
                 desc, 14, LIGHT_GRAY, False)

# ═══════════════════ SLIDE 5: ARCHITECTURE OVERVIEW ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "SYSTEM ARCHITECTURE", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
             "Two-Phase Pipeline: Self-Supervised Pretraining + Supervised Fine-Tuning",
             18, YELLOW, True, PP_ALIGN.LEFT)

# Visual stream
add_rect(slide, Inches(0.8), Inches(2.4), Inches(2.2), Inches(0.8), DARK_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(2.45), Inches(2), Inches(0.35),
             "Video Input", 14, ACCENT_CYAN, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.9), Inches(2.8), Inches(2), Inches(0.3),
             "16 frames, 224x224", 12, MUTED_GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(3.1), Inches(2.6), Inches(0.5), Inches(0.4),
             "\u2192", 24, ACCENT_BLUE, True, PP_ALIGN.CENTER)

add_rect(slide, Inches(3.6), Inches(2.4), Inches(2.4), Inches(0.8), DARK_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(3.7), Inches(2.45), Inches(2.2), Inches(0.35),
             "ViT-B/16 Encoder", 14, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.7), Inches(2.8), Inches(2.2), Inches(0.3),
             "(B, 16, 768)", 12, MUTED_GRAY, False, PP_ALIGN.CENTER)

# Audio stream
add_rect(slide, Inches(0.8), Inches(3.7), Inches(2.2), Inches(0.8), DARK_CARD, ORANGE)
add_text_box(slide, Inches(0.9), Inches(3.75), Inches(2), Inches(0.35),
             "Audio Input", 14, ORANGE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.9), Inches(4.1), Inches(2), Inches(0.3),
             "128-bin Mel Spectrogram", 12, MUTED_GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(3.1), Inches(3.9), Inches(0.5), Inches(0.4),
             "\u2192", 24, ORANGE, True, PP_ALIGN.CENTER)

add_rect(slide, Inches(3.6), Inches(3.7), Inches(2.4), Inches(0.8), DARK_CARD, ORANGE)
add_text_box(slide, Inches(3.7), Inches(3.75), Inches(2.2), Inches(0.35),
             "CNN Encoder (4L)", 14, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.7), Inches(4.1), Inches(2.2), Inches(0.3),
             "(B, T_a, 768)", 12, MUTED_GRAY, False, PP_ALIGN.CENTER)

# Arrows to fusion
add_text_box(slide, Inches(6.1), Inches(2.6), Inches(0.7), Inches(0.4),
             "\u2192", 24, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.1), Inches(3.9), Inches(0.7), Inches(0.4),
             "\u2192", 24, ORANGE, True, PP_ALIGN.CENTER)

# Fusion Engine
add_rect(slide, Inches(6.8), Inches(2.4), Inches(3), Inches(2.1), DARK_CARD, GREEN)
add_text_box(slide, Inches(6.9), Inches(2.5), Inches(2.8), Inches(0.4),
             "Multimodal Bottleneck", 16, GREEN, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.9), Inches(2.9), Inches(2.8), Inches(0.4),
             "Transformer (MBT)", 16, GREEN, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.9), Inches(3.4), Inches(2.8), Inches(0.3),
             "4 layers | 4 bottleneck tokens", 12, MUTED_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.9), Inches(3.7), Inches(2.8), Inches(0.3),
             "8 attention heads | 0.1 dropout", 12, MUTED_GRAY, False, PP_ALIGN.CENTER)

# Arrow to classifier
add_text_box(slide, Inches(9.9), Inches(3.2), Inches(0.7), Inches(0.4),
             "\u2192", 24, GREEN, True, PP_ALIGN.CENTER)

# Classifier
add_rect(slide, Inches(10.5), Inches(2.7), Inches(2.4), Inches(1.5), DARK_CARD, YELLOW)
add_text_box(slide, Inches(10.6), Inches(2.8), Inches(2.2), Inches(0.4),
             "MLP Classifier", 16, YELLOW, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.6), Inches(3.2), Inches(2.2), Inches(0.3),
             "1536 \u2192 512 \u2192 256 \u2192 1", 12, MUTED_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.6), Inches(3.6), Inches(2.2), Inches(0.3),
             "GELU + Dropout", 12, MUTED_GRAY, False, PP_ALIGN.CENTER)

# Component summary table
comp_data = [
    ["Component", "Description", "Output"],
    ["Visual Encoder", "ViT-B/16 + Temporal Pos. Embeddings", "(B, 16, 768)"],
    ["Audio Encoder", "4-layer CNN \u2192 768-dim projection", "(B, T_a, 768)"],
    ["Fusion Engine", "4-layer MBT, 4 bottleneck tokens", "(B, 768) per mod."],
    ["Classifier", "3-layer MLP (GELU, Dropout)", "Binary logit"],
]
add_table_slide_content(slide, Inches(0.8), Inches(5.2), Inches(11.7), comp_data,
                        [Inches(2.5), Inches(5.5), Inches(3.7)])

# ═══════════════════ SLIDE 6: SSL PRETRAINING ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "PHASE 1: SELF-SUPERVISED PRETRAINING", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(4.5), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "Pretrained on VoxCeleb (real celebrity videos) \u2014 No fake labels required",
             18, YELLOW, True)

# Three SSL objectives
ssl_tasks = [
    ("Contrastive Learning", ACCENT_BLUE,
     ["MoCo-style momentum encoder (m=0.999)",
      "2048-sample memory bank",
      "Temperature \u03c4 = 0.07",
      "InfoNCE loss, weight = 0.5",
      "Variance regularization (w=0.1)"]),
    ("AV Sync Prediction", GREEN,
     ["Binary: synced vs desynced",
      "50% synced, 50% desynced",
      "30% cross-sample audio swap",
      "70% temporal shift (15\u201340 frames)",
      "Cross-entropy loss, weight = 1.0"]),
    ("Masked Reconstruction", ORANGE,
     ["Visual: 85% mask, block size 8",
      "Audio: 80% mask, block size 12",
      "Smooth L1 reconstruction loss",
      "Audio weight = 2.0",
      "Visual weight = 3.0"]),
]

for i, (title, color, bullets) in enumerate(ssl_tasks):
    x = Inches(0.8) + Inches(i * 4.1)
    add_rect(slide, x, Inches(2.3), Inches(3.8), Inches(3.4), DARK_CARD, color)
    add_text_box(slide, x + Inches(0.15), Inches(2.4), Inches(3.5), Inches(0.4),
                 title, 17, color, True, PP_ALIGN.CENTER)
    for j, bullet in enumerate(bullets):
        add_text_box(slide, x + Inches(0.2), Inches(3.0) + Inches(j * 0.45), Inches(3.4), Inches(0.4),
                     f"\u2022  {bullet}", 13, LIGHT_GRAY, False)

# Training config at bottom
add_text_box(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.5),
             "Training: 50 epochs | Batch 8 (x4 accumulation = 32 effective) | AdamW lr=1e-4 | Mixed Precision (FP16) | Grad clip 1.0",
             14, MUTED_GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════ SLIDE 7: FINE-TUNING ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "PHASE 2: SUPERVISED FINE-TUNING", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(4), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
             "Fine-tuned on FakeAVCeleb for binary classification (Real vs Fake)",
             18, YELLOW, True)

# Details
ft_details = [
    ("Dataset", "FakeAVCeleb \u2014 real and fake videos with RARA/FAFA/FARA/RAFA categories"),
    ("Split", "80% training / 20% validation"),
    ("Loss Function", "Binary Cross-Entropy with Logits (BCE)"),
    ("Scheduler", "Cosine Annealing (T_max = 20 epochs)"),
    ("Best Model", "Saved by best validation F1 score"),
]
for i, (label, value) in enumerate(ft_details):
    y = Inches(2.3) + Inches(i * 0.55)
    add_text_box(slide, Inches(1.2), y, Inches(2.5), Inches(0.45),
                 label, 16, ACCENT_CYAN, True)
    add_text_box(slide, Inches(3.8), y, Inches(8.5), Inches(0.45),
                 value, 16, LIGHT_GRAY, False)

# Differential LR highlight
add_rect(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5), DARK_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(5.3), Inches(11), Inches(0.4),
             "Differential Learning Rates (Key Strategy)", 18, ACCENT_BLUE, True)
add_text_box(slide, Inches(1.4), Inches(5.8), Inches(10.5), Inches(0.4),
             "\u2022  Pretrained encoders & fusion: lr = 1e-5 (preserves SSL knowledge)", 15, LIGHT_GRAY, False)
add_text_box(slide, Inches(1.4), Inches(6.2), Inches(10.5), Inches(0.4),
             "\u2022  New classifier head: lr = 1e-4 (learns decision boundary quickly)", 15, LIGHT_GRAY, False)

# ═══════════════════ SLIDE 8: RESULTS ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "RESULTS & METRICS", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3), ACCENT_BLUE)

# SSL Convergence table
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
             "SSL Pretraining Convergence", 20, YELLOW, True)
ssl_data = [
    ["Loss Component", "Initial", "Final (15 ep.)"],
    ["Total Loss", "3.60", "0.55"],
    ["Contrastive Loss", "4.50", "0.40"],
    ["Sync Loss", "0.65", "0.05"],
    ["Audio Recon Loss", "0.40", "0.20"],
    ["Visual Recon Loss", "0.25", "0.05"],
]
add_table_slide_content(slide, Inches(0.8), Inches(2.1), Inches(5.5), ssl_data,
                        [Inches(2.2), Inches(1.65), Inches(1.65)], GREEN)

# Fine-tuning metrics (big numbers)
add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
             "Fine-Tuning Performance (FakeAVCeleb)", 20, YELLOW, True)

metrics = [
    ("85.2%", "Accuracy", ACCENT_BLUE),
    ("86.42%", "Precision", GREEN),
    ("75.6%", "Recall", ORANGE),
    ("80.1%", "F1 Score", ACCENT_CYAN),
]
for i, (value, label, color) in enumerate(metrics):
    col = i % 2
    row = i // 2
    x = Inches(7) + Inches(col * 2.8)
    y = Inches(2.2) + Inches(row * 1.6)
    add_rect(slide, x, y, Inches(2.5), Inches(1.35), DARK_CARD, color)
    add_text_box(slide, x, y + Pt(8), Inches(2.5), Inches(0.6),
                 value, 32, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Pt(50), Inches(2.5), Inches(0.4),
                 label, 14, MUTED_GRAY, False, PP_ALIGN.CENTER)

# Observations
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.4),
             "Key Observations", 18, ACCENT_BLUE, True)
observations = [
    "All SSL loss components converged steadily \u2014 model learned meaningful cross-modal representations",
    "High precision (86.42%) means low false accusation rate \u2014 critical for real-world deployment",
    "Precision > Recall gap indicates conservative model; threshold tuning can adjust this trade-off",
    "Multimodal approach outperforms visual-only methods by capturing audio-visual inconsistencies",
]
for i, obs in enumerate(observations):
    add_text_box(slide, Inches(1.2), Inches(6.0) + Inches(i * 0.36), Inches(11.3), Inches(0.35),
                 f"\u2022  {obs}", 13, LIGHT_GRAY, False)

# ═══════════════════ SLIDE 9: TECH STACK ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "TECHNOLOGY STACK", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3), ACCENT_BLUE)

tech_data = [
    ["Layer", "Technology", "Role"],
    ["Language", "Python 3.8+", "All scripts, model definitions, training"],
    ["DL Framework", "PyTorch 2.0+", "Neural networks, autograd, CUDA"],
    ["Vision", "TorchVision", "ViT-B/16 backbone, video I/O, transforms"],
    ["Audio", "TorchAudio", "MelSpectrogram, AmplitudeToDB, Resample"],
    ["Monitoring", "TensorBoard", "Real-time loss tracking during SSL"],
    ["Visualization", "Matplotlib", "Loss curves, metrics dashboard plots"],
    ["Progress", "tqdm", "Progress bars for training loops"],
    ["Data Download", "gdown, requests", "FakeAVCeleb (Drive), VoxCeleb (HTTP)"],
]
add_table_slide_content(slide, Inches(0.8), Inches(1.6), Inches(11.7), tech_data,
                        [Inches(2.2), Inches(3.5), Inches(6)])

# Hardware
add_text_box(slide, Inches(0.8), Inches(5.8), Inches(5), Inches(0.4),
             "Hardware Requirements", 18, YELLOW, True)
hw_data = [
    ["Component", "Minimum", "Recommended"],
    ["GPU", "NVIDIA T4 (16 GB)", "RTX 3090/4090 (24 GB)"],
    ["RAM", "16 GB", "32 GB"],
    ["Storage", "50 GB", "100 GB"],
]
add_table_slide_content(slide, Inches(0.8), Inches(6.3), Inches(7), hw_data,
                        [Inches(2), Inches(2.5), Inches(2.5)], ORANGE)

# ═══════════════════ SLIDE 10: DATA PIPELINE ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "DATA PIPELINE & PREPROCESSING", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(4), ACCENT_BLUE)

# Visual pipeline
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4),
             "Visual Preprocessing", 20, ACCENT_BLUE, True)
vis_steps = [
    ("1. Temporal Sampling", "16 frames uniformly sampled via linspace"),
    ("2. Center Crop", "80% of min spatial dimension"),
    ("3. Resize", "224 x 224 pixels"),
    ("4. Normalize", "ImageNet mean & std"),
]
for i, (step, detail) in enumerate(vis_steps):
    y = Inches(2.1) + Inches(i * 0.55)
    add_rect(slide, Inches(0.8), y, Inches(5.5), Inches(0.45), DARK_CARD)
    add_text_box(slide, Inches(1), y + Pt(2), Inches(2.2), Inches(0.4),
                 step, 14, ACCENT_CYAN, True)
    add_text_box(slide, Inches(3.2), y + Pt(2), Inches(3), Inches(0.4),
                 detail, 14, LIGHT_GRAY, False)

# Audio pipeline
add_text_box(slide, Inches(7), Inches(1.6), Inches(5), Inches(0.4),
             "Audio Preprocessing", 20, ORANGE, True)
aud_steps = [
    ("1. Resample", "Convert to 16 kHz mono"),
    ("2. Mel Spectrogram", "128 bins, 1024 FFT, 512 hop"),
    ("3. Log Scale", "AmplitudeToDB (power, 80 dB)"),
    ("4. Fix Length", "Pad/trim to 94 time steps"),
]
for i, (step, detail) in enumerate(aud_steps):
    y = Inches(2.1) + Inches(i * 0.55)
    add_rect(slide, Inches(7), y, Inches(5.5), Inches(0.45), DARK_CARD)
    add_text_box(slide, Inches(7.2), y + Pt(2), Inches(2.2), Inches(0.4),
                 step, 14, ORANGE, True)
    add_text_box(slide, Inches(9.4), y + Pt(2), Inches(3), Inches(0.4),
                 detail, 14, LIGHT_GRAY, False)

# Datasets section
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(10), Inches(0.4),
             "Datasets Used", 20, YELLOW, True)
ds_data = [
    ["Dataset", "Purpose", "Content"],
    ["VoxCeleb", "SSL Pretraining (Phase 1)", "Real celebrity interview videos (no fakes)"],
    ["FakeAVCeleb", "Fine-Tuning (Phase 2)", "Real + Fake videos (RARA, FAFA, FARA, RAFA)"],
]
add_table_slide_content(slide, Inches(0.8), Inches(5.1), Inches(11.7), ds_data,
                        [Inches(2.5), Inches(3.5), Inches(5.7)])

# Augmentations
add_text_box(slide, Inches(0.8), Inches(6.1), Inches(5), Inches(0.4),
             "SSL Data Augmentation", 18, ACCENT_BLUE, True)
aug_items = [
    "Video: Horizontal flip (50%), Color jitter (80%), Gaussian blur (10%)",
    "Audio: Gain adjust (50%), Time masking (2 blocks x10), Freq masking (2 blocks x15)",
]
for i, item in enumerate(aug_items):
    add_text_box(slide, Inches(1.2), Inches(6.5) + Inches(i * 0.4), Inches(11.5), Inches(0.35),
                 f"\u2022  {item}", 13, LIGHT_GRAY, False)

# ═══════════════════ SLIDE 11: PROJECT STRUCTURE ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "PROJECT STRUCTURE", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3), ACCENT_BLUE)

# Models folder
sections = [
    ("Models/", ACCENT_BLUE, [
        ("Visual_encoder.py", "ViT-based video encoder with temporal embeddings"),
        ("audio_encoder.py", "4-layer CNN audio tokenizer"),
        ("Fusion_Engine.py", "Multimodal Bottleneck Transformer (MBT)"),
        ("ssl_tasks.py", "SSL pretraining orchestrator with 3 objectives"),
    ]),
    ("utils/", GREEN, [
        ("visual_preprocessing.py", "Frame sampling, cropping, normalization"),
        ("audio_preprocessing.py", "Mel spectrogram extraction"),
        ("augmentations.py", "Video and audio augmentation strategies"),
    ]),
    ("dataset/", ORANGE, [
        ("pretrain_datset_download.py", "VoxCeleb dataset downloader"),
        ("preprocess_dataset.py", "Video \u2192 tensor offline conversion"),
        ("download_finetune_dataset.py", "FakeAVCeleb dataset downloader"),
        ("prepare_finetune_dataset.py", "Real/fake directory organizer"),
    ]),
    ("Root Scripts", YELLOW, [
        ("ssl_pretrain.py", "Self-supervised pretraining script (Phase 1)"),
        ("finetune_deepfake.py", "Supervised fine-tuning script (Phase 2)"),
    ]),
]

y_offset = Inches(1.6)
for sec_name, sec_color, files in sections:
    add_text_box(slide, Inches(0.8), y_offset, Inches(3), Inches(0.35),
                 sec_name, 16, sec_color, True)
    for fname, desc in files:
        y_offset += Inches(0.32)
        add_text_box(slide, Inches(1.5), y_offset, Inches(3.5), Inches(0.3),
                     fname, 12, WHITE, False)
        add_text_box(slide, Inches(5.2), y_offset, Inches(7.5), Inches(0.3),
                     desc, 12, MUTED_GRAY, False)
    y_offset += Inches(0.45)

# ═══════════════════ SLIDE 12: LIMITATIONS & FUTURE ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "LIMITATIONS & FUTURE SCOPE", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(4), ACCENT_BLUE)

# Limitations
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
             "Current Limitations", 20, RED_ACCENT, True)
limitations = [
    "GPU with 16+ GB VRAM required for training",
    "VoxCeleb/FakeAVCeleb may not cover all\n  deepfake generation techniques",
    "Requires both audio and video; no fallback\n  for audio-missing scenarios",
    "Batch processing only; not real-time ready",
    "Not tested against adversarial attacks",
]
for i, lim in enumerate(limitations):
    y = Inches(2.0) + Inches(i * 0.55)
    add_text_box(slide, Inches(1.2), y, Inches(5.2), Inches(0.5),
                 f"\u2022  {lim}", 13, LIGHT_GRAY, False)

# Future scope
add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
             "Future Scope", 20, GREEN, True)
future = [
    "Web application (FastAPI/Flask) for browser-\n  based video upload & detection",
    "Real-time video analysis via quantization,\n  distillation, or TensorRT",
    "Additional modalities: optical flow, facial\n  landmarks, frequency-domain analysis",
    "Adversarial training for robustness",
    "Cross-dataset generalization testing\n  (FaceForensics++, Celeb-DF, DFDC)",
    "Explainability: attention & saliency maps",
    "Model compression for mobile deployment",
]
for i, item in enumerate(future):
    y = Inches(2.0) + Inches(i * 0.55)
    add_text_box(slide, Inches(7.2), y, Inches(5.2), Inches(0.5),
                 f"\u2022  {item}", 13, LIGHT_GRAY, False)

# ═══════════════════ SLIDE 13: CONCLUSION ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "CONCLUSION", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_BLUE)

conclusions = [
    ("Multimodal approach is effective",
     "Combining audio and visual analysis through the Multimodal Bottleneck Transformer captures cross-modal inconsistencies that single-modality methods miss."),
    ("Self-supervised pretraining works",
     "Pretraining on unlabelled VoxCeleb with contrastive, sync, and masked objectives provides a strong foundation that transfers well to deepfake classification."),
    ("Competitive results achieved",
     "85.2% accuracy, 86.42% precision, 75.6% recall, and 80.1% F1 on FakeAVCeleb, with high precision minimizing false accusations."),
    ("Modular and extensible design",
     "Each component (ViT encoder, CNN encoder, MBT fusion, MLP classifier) can be independently upgraded as better techniques emerge."),
    ("Practical and efficient",
     "Mixed-precision training, gradient accumulation, and differential learning rates make the system trainable on consumer GPUs (T4+)."),
]

for i, (title, desc) in enumerate(conclusions):
    y = Inches(1.6) + Inches(i * 1.1)
    add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.9), DARK_CARD, GREEN)
    add_text_box(slide, Inches(1.1), y + Pt(5), Inches(11), Inches(0.35),
                 f"\u2713  {title}", 17, GREEN, True)
    add_text_box(slide, Inches(1.5), y + Pt(28), Inches(10.5), Inches(0.45),
                 desc, 14, LIGHT_GRAY, False)

# ═══════════════════ SLIDE 14: REFERENCES ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "REFERENCES", 32, ACCENT_CYAN, True)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2.5), ACCENT_BLUE)

references = [
    "[1] Dosovitskiy et al. (2021) \u2014 An Image is Worth 16x16 Words: Transformers for Image Recognition at Scale. ICLR.",
    "[2] Nagrani et al. (2021) \u2014 Attention Bottlenecks for Multimodal Fusion. NeurIPS.",
    "[3] He et al. (2020) \u2014 Momentum Contrast for Unsupervised Visual Representation Learning. CVPR.",
    "[4] Khalid et al. (2021) \u2014 FakeAVCeleb: A Novel Audio-Video Multimodal Deepfake Dataset. NeurIPS.",
    "[5] Nagrani et al. (2017) \u2014 VoxCeleb: A Large-Scale Speaker Identification Dataset. Interspeech.",
    "[6] Park et al. (2019) \u2014 SpecAugment: A Simple Data Augmentation Method for ASR. Interspeech.",
    "[7] Loshchilov & Hutter (2019) \u2014 Decoupled Weight Decay Regularization. ICLR.",
    "[8] Micikevicius et al. (2018) \u2014 Mixed Precision Training. ICLR.",
    "[9] Touvron et al. (2021) \u2014 Training Data-Efficient Image Transformers. ICML.",
    "[10] Zi et al. (2020) \u2014 WildDeepfake: A Challenging Real-World Dataset. ACM MM.",
]
for i, ref in enumerate(references):
    add_text_box(slide, Inches(1), Inches(1.6) + Inches(i * 0.52), Inches(11.5), Inches(0.45),
                 ref, 14, LIGHT_GRAY, False)

# ═══════════════════ SLIDE 15: THANK YOU ═══════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(1.5), Inches(2), Inches(10), Inches(1),
             "THANK YOU", 48, ACCENT_CYAN, True, PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5), Inches(3.2), Inches(3.333), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.6),
             "Questions & Discussion", 24, YELLOW, True, PP_ALIGN.CENTER)

add_multiline_text(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(2), [
    ("Sankalp (12311001)  |  Madhur Suman (12311002)", 18, WHITE, True),
    ("", 10, WHITE, False),
    ("Supervised by: Dr. Md. Arquam", 16, MUTED_GRAY, False),
    ("Indian Institute of Information Technology, Sonepat", 16, MUTED_GRAY, False),
    ("April 2026", 16, MUTED_GRAY, False),
])
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        for p in shape.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER

# ─── SAVE ──────────────────────────────────────────────────
output_path = "DeepDetect2_Minor_Project_Evaluation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
